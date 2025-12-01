"""Document converter implementation for text-based files."""

from __future__ import annotations

import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, List, Optional
# 处理相对导入，支持直接运行脚本
current_dir = Path(__file__).parent
parent_dir = current_dir.parent.parent
if str(parent_dir) not in sys.path:
    sys.path.insert(0, str(parent_dir))
from multimodal.converter import ConversionChunk, ConversionOutput, MultimodalConverter
from multimodal.factory import ConverterFactory

class DocumentConverter(MultimodalConverter):
    SUPPORTED_EXTENSIONS: List[str] = [
        "pdf",
        "docx",
        "doc",
        "md",
        "txt",
        "pptx",
    ]

    def convert(self, source, *, source_type: str = "file_path", **kwargs: Any) -> ConversionOutput:
        """
        处理文档文件，将内容转换为文本 chunks
        
        处理流程：
        1. 读取文件内容（根据文件类型使用不同的解析方法）
        2. 如果文本太长，使用 _chunk_text() 分割成多个 chunks
        3. 返回 ConversionOutput
        """
        try:
            if source_type != "file_path":
                raise ValueError(f"DocumentConverter 只支持本地文件路径 (source_type='file_path')，当前为: {source_type}")
            
            file_path = Path(source)
            
            # 验证文件是否存在
            if not file_path.exists():
                raise FileNotFoundError(f"文件不存在: {file_path}")
            
            file_ext = file_path.suffix.lower().lstrip('.')
            
            self._report_progress(0.1, f"正在读取文件: {file_path.name}...")
            
            # 根据文件类型读取内容
            text_content = self._read_file_content(file_path, file_ext)
            
            if not text_content:
                return ConversionOutput(
                    status="failed",
                    chunks=[],
                    metadata={"converter_provider": "document_converter"},
                    error=f"无法读取文件内容: {file_path}",
                )
            
            self._report_progress(0.5, f"文件读取完成，内容长度: {len(text_content)} 字符")
            
            # 如果文本太长，分割成多个 chunks
            # 使用基类提供的 _chunk_text 方法
            if len(text_content) > self.max_chunk_tokens * 4:  # 粗略估算：每个token约4个字符
                self._report_progress(0.7, "文本较长，正在分割成多个 chunks...")
                chunks = self._chunk_text(
                    text_content,
                    chunk_size=self.max_chunk_tokens,
                    overlap=200,
                    base_metadata={"source_type": "document", "file_extension": file_ext},
                )
            else:
                # 文本不长，直接作为一个 chunk
                chunks = [
                    ConversionChunk(
                        text=text_content,
                        chunk_index=0,
                        metadata={
                            "source_type": "document",
                            "file_extension": file_ext,
                            "chunk_count_estimate": 1,
                        },
                    )
                ]
            
            self._report_progress(1.0, f"文档转换完成，共 {len(chunks)} 个 chunks")
            
            return ConversionOutput(
                status="success",
                chunks=chunks,
                metadata={
                    "converter_provider": "document_converter",
                    "converter_version": "1.0.0",
                    "conversion_time": datetime.utcnow().isoformat() + "Z",
                    "file_extension": file_ext,
                    "total_chunks": len(chunks),
                },
            )
        except Exception as e:
            return ConversionOutput(
                status="failed",
                chunks=[],
                metadata={"converter_provider": "document_converter"},
                error=str(e),
            )

    def _read_file_content(self, file_path: Path, file_ext: str) -> str:
        """
        根据文件类型读取内容
        目前支持：txt, md（直接读取）
        其他格式（pdf, docx, pptx）需要相应的库解析
        """
        # 简单文本文件：直接读取
        if file_ext in ["txt", "md"]:
            try:
                # 尝试多种编码
                encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            return f.read()
                    except UnicodeDecodeError:
                        continue
                # 如果所有编码都失败，使用二进制模式读取
                with open(file_path, 'rb') as f:
                    return f.read().decode('utf-8', errors='ignore')
            except Exception as e:
                raise ValueError(f"读取文件失败: {e}")
        
        # PDF 文件
        elif file_ext == "pdf":
            try:
                import PyPDF2  # type: ignore
                text_parts = []
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        text_parts.append(page.extract_text())
                return '\n\n'.join(text_parts)
            except ImportError:
                raise ImportError("处理 PDF 文件需要安装 PyPDF2: pip install PyPDF2")
            except Exception as e:
                raise ValueError(f"解析 PDF 文件失败: {e}")
        
        # DOCX 文件
        elif file_ext == "docx":
            try:
                from docx import Document  # type: ignore
                doc = Document(file_path)
                paragraphs = [para.text for para in doc.paragraphs]
                return '\n'.join(paragraphs)
            except ImportError:
                raise ImportError("处理 DOCX 文件需要安装 python-docx: pip install python-docx")
            except Exception as e:
                raise ValueError(f"解析 DOCX 文件失败: {e}")
        
        # PPTX 文件
        elif file_ext == "pptx":
            try:
                from pptx import Presentation  # type: ignore
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                return '\n\n'.join(text_parts)
            except ImportError:
                raise ImportError("处理 PPTX 文件需要安装 python-pptx: pip install python-pptx")
            except Exception as e:
                raise ValueError(f"解析 PPTX 文件失败: {e}")
        
        # DOC 文件（旧版 Word 格式，需要特殊处理）
        elif file_ext == "doc":
            raise NotImplementedError("DOC 文件（旧版 Word 格式）暂不支持，请转换为 DOCX 格式")
        
        else:
            raise ValueError(f"不支持的文件类型: {file_ext}")

    def supports(self, *, file_type: str, mime_type: str = None) -> bool:
        return file_type.lower() in self.SUPPORTED_EXTENSIONS


# Register the converter
ConverterFactory.register("document", DocumentConverter, priority=0)


def main():
    """测试 DocumentConverter 功能"""
    import sys
    
    # 检查是否提供了文件路径
    if len(sys.argv) < 2:
        print("用法: python file_converter.py <文件路径>")
        print("示例: python file_converter.py document.txt")
        print("示例: python file_converter.py document.pdf")
        print("示例: python file_converter.py document.docx")
        print()
        print("支持的文件类型:")
        print("  - txt, md: 直接读取文本")
        print("  - pdf: 需要安装 PyPDF2 (pip install PyPDF2)")
        print("  - docx: 需要安装 python-docx (pip install python-docx)")
        print("  - pptx: 需要安装 python-pptx (pip install python-pptx)")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    # 检查文件是否存在
    if not os.path.exists(file_path):
        print(f"错误: 文件不存在: {file_path}")
        sys.exit(1)
    
    print("=" * 60)
    print("DocumentConverter 测试")
    print("=" * 60)
    print(f"文件路径: {file_path}")
    print(f"文件类型: {Path(file_path).suffix}")
    print()
    
    # 进度回调函数
    def progress_callback(progress: float, message: str):
        print(f"[进度 {progress*100:.1f}%] {message}")
    
    try:
        # 创建转换器实例
        converter = DocumentConverter(
            progress_callback=progress_callback,
        )
        
        print("开始转换文档...")
        print()
        
        # 执行转换
        result = converter.convert(file_path, source_type="file_path")
        
        print()
        print("=" * 60)
        print("转换结果")
        print("=" * 60)
        print(f"状态: {result.status}")
        
        if result.status == "success":
            print(f"Chunk 数量: {len(result.chunks)}")
            print()
            
            for i, chunk in enumerate(result.chunks):
                print(f"--- Chunk {i+1}/{len(result.chunks)} ---")
                print(f"文本长度: {len(chunk.text)} 字符")
                print(f"元数据: {chunk.metadata}")
                print()
                
                # 只显示前500个字符，避免输出过长
                preview_text = chunk.text[:500]
                if len(chunk.text) > 500:
                    preview_text += f"\n... (还有 {len(chunk.text) - 500} 个字符)"
                
                print("文档内容预览:")
                print("-" * 60)
                print(preview_text)
                print("-" * 60)
                print()
            
            print("输出元数据:")
            print(result.metadata)
        else:
            print(f"错误: {result.error}")
            if result.error:
                print(f"错误详情: {result.error}")
        
    except ImportError as e:
        print(f"导入错误: {e}")
        print()
        print("请安装相应的依赖库:")
        file_ext = Path(file_path).suffix.lower().lstrip('.')
        if file_ext == "pdf":
            print("  pip install PyPDF2")
        elif file_ext == "docx":
            print("  pip install python-docx")
        elif file_ext == "pptx":
            print("  pip install python-pptx")
        sys.exit(1)
    except Exception as e:
        print(f"转换失败: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
if __name__ == "__main__":
    main()